#!/usr/bin/env python3
"""build-overview-pptx.py · 《WorkLoom 获客系统 · 业务全景介绍》PPT 生成器（给人类开发者）
视角：业务 + 客户。产出 docs/system-overview.pptx（16:9，糖果系，13 页）。
用法：python3 scripts/build-overview-pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ROOT = os.path.join(os.path.dirname(__file__), "..")
SHOTS = os.path.join(ROOT, "docs/demo/preview-shots")

BG = RGBColor(0xFF, 0xF5, 0xF7); MAIN = RGBColor(0xD4, 0x00, 2 * 0x15)  # 蜜桃雾 / 深珊瑚
MAIN = RGBColor(0xD4, 0x00, 0x2A)
TEXT = RGBColor(0x4A, 0x2B, 0x33); SUB = RGBColor(0x8A, 0x4B, 0x5A); WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.33), Inches(7.5)
BLANK = prs.slide_layouts[6]

def slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = MAIN if dark else BG
    return s

def text(s, x, y, w, h, lines):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (t, size, color, bold, *rest) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = rest[0] if rest else PP_ALIGN.LEFT
        r = p.add_run(); r.text = t
        f = r.font; f.size = Pt(size); f.color.rgb = color; f.bold = bold; f.name = "PingFang SC"
        p.space_after = Pt(4)
    return tb

def bullets(s, x, y, w, h, items, size=15, gap=6):
    lines = []
    for head, body in items:
        lines.append((head, size, MAIN, True))
        if body: lines.append((body, size - 3, SUB, False))
    return text(s, x, y, w, h, lines)

# ① 封面
s = slide(dark=True)
text(s, 1, 2.0, 11.3, 3.6, [
    ("WorkLoom 获客系统", 48, WHITE, True, PP_ALIGN.CENTER),
    ("短视频社媒营销 × GEO × 获客五环", 22, WHITE, False, PP_ALIGN.CENTER),
    ("把获客从玄学，变成一门可测量、可追溯、可归因的生意", 18, WHITE, False, PP_ALIGN.CENTER),
    ("—— 给开发者的业务全景介绍 ——", 14, RGBColor(0xFF, 0xD6, 0xDE), False, PP_ALIGN.CENTER),
])

# ② 客户的痛
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("你的客户，正在被这三笔账拖垮", 30, MAIN, True)])
bullets(s, 0.9, 1.5, 11.5, 4.2, [
    ("第一笔：流量越来越贵", "投流成本年年涨，内容不发就没声量，发了没转化就是给平台打工"),
    ("第二笔：询盘接不住", "评论区一句「怎么订」半小时没人回，客户就去别家了——夜间/高峰期询盘流失率超 70%"),
    ("第三笔：佣金被抽走", "以酒店为例 OTA 佣金 15-25%——辛苦服务一整晚，平台躺着抽走四分之一"),
    ("还有一个大多数人没看见的变化", "42% 消费者已开始用生成式 AI 查消费信息——AI 搜索里没有你，就失去了一个爆发中的新入口"),
], size=18)

# ③ 系统是什么
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("WorkLoom 获客系统是什么", 30, MAIN, True)])
text(s, 0.9, 1.4, 11.5, 1.2, [
    ("一支 AI 获客班组住进商家的通讯录：找人群、做内容、接询盘、跟线索、算成交。", 19, TEXT, True),
    ("商家只做三件事：定方向、拍板、收钱。", 19, MAIN, True),
])
bullets(s, 0.9, 2.9, 11.5, 3.8, [
    ("不是代运营工具，是「经营系统」", "员工、围栏、审批、账本、事件链一应俱全——AI 在规则内干活，人在关键处拍板"),
    ("不是 demo 玩具", "多租户底座、RLS 数据隔离、400+ 条场景测试、发布门禁——按生产标准建造"),
    ("首垂直行业：酒店", "低星单体/民宿/无人酒店三类客群默认配置；底座按 Bundle 复制到各行业"),
], size=17)

# ④ 客户视角的一天
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("客户视角：酒店店主的一天", 30, MAIN, True)])
rows = [
    ("07:30", "手机收到【晨报】：昨日曝光/询盘/成交、待拍板 2 项，一句话看懂生意"),
    ("09:15", "AI 已把昨夜评论区的「怎么订」全部秒回，3 条高意向线索转成了企微私聊"),
    ("11:00", "内容工厂按本周热点出了 5 条短视频脚本，店主点「批准」2 条进入发布队列"),
    ("15:20", "RPA 模拟人工把视频发到抖音/小红书，回执自动归档，失败挂起转人工"),
    ("21:40", "GEO 监测：「杭州亲子酒店」AI 搜索答案里出现了自家店，曝光位次上升"),
    ("23:55", "夜班接管：继续回询盘、发券跟进线索；店主早已休息"),
]
for i, (t, d) in enumerate(rows):
    text(s, 0.9, 1.5 + i * 0.95, 11.5, 0.9, [(f"{t}  {d}", 15, TEXT, False)])

# ⑤ 获客五环
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("获客五环：每一环都可度量到钱", 30, MAIN, True)])
rings = [
    ("① 意图洞察", "社媒监听 + GEO 情报：客户在问什么、AI 搜索缺什么", "知道该做什么内容"),
    ("② 双域触达", "短视频内容工厂 + GEO 内容工厂，一条生产线两个战场", "内容不再断更"),
    ("③ 四路承接", "评论秒回 / 私信 / 企微 / C 端服务前台，询盘不漏一个", "询盘流失率从 70% 压到个位数"),
    ("④ 线索转化", "券 SKU + 私域跟进 + 人审报价，绕开 OTA 高佣金", "佣金节省看得见（演示口径月省 ¥14,900）"),
    ("⑤ 归因复盘", "五元事件全链路归因：每块钱从哪来、到哪去", "下月预算花在刀刃上"),
]
for i, (a, b, c) in enumerate(rings):
    y = 1.5 + i * 1.12
    text(s, 0.9, y, 3.2, 1.0, [(a, 17, MAIN, True)])
    text(s, 4.1, y, 5.6, 1.0, [(b, 13, TEXT, False)])
    text(s, 9.9, y, 3.0, 1.0, [(c, 13, SUB, False)])

# ⑥ 三端界面（截图）
s = slide()
text(s, 0.8, 0.35, 11.7, 0.8, [("三端界面：客户每天真正在用的样子", 28, MAIN, True)])
imgs = [("pc-3000.png", "PC · B 端工作台"), ("shell-guest.png", "B 端移动（手机壳）"), ("mobile-c-3002.png", "C 端 AI 服务前台")]
for i, (f, cap) in enumerate(imgs):
    p = os.path.join(SHOTS, f)
    x = 0.7 + i * 4.25
    if os.path.exists(p):
        s.shapes.add_picture(p, Inches(x), Inches(1.3), width=Inches(4.0))
    text(s, x, 3.95, 4.0, 0.5, [(cap, 14, MAIN, True, PP_ALIGN.CENTER)])
text(s, 0.9, 6.7, 11.5, 0.6, [("全部为运行态实机截图（Mock 数据模式），pnpm preview:all 一键复现", 13, SUB, False, PP_ALIGN.CENTER)])

# ⑦ AI 班组
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("AI 班组：16 个员工各管一摊，7×24 不下班", 30, MAIN, True)])
bullets(s, 0.9, 1.5, 11.5, 4.8, [
    ("内容线", "热点猎手找选题 → 脚本写手出稿 → 人审台把关 → RPA 六平台分发 → 数据回收复盘"),
    ("营销线", "意图雷达盯评论区与 AI 搜索 → GEO 写手抢占 AI 答案位 → 曝光监测每日报位次"),
    ("服务线", "AI 接待员秒回询盘（知识库 385 问）→ 高意向转私域 → 券运营发券促单 → 客成回访"),
    ("经营线", "公司 CEO 统筹晨报 → 收益看守盯价格库存 → 巡检发现异常 → 夜班接管不打烊"),
], size=17)

# ⑧ 信任设计
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("客户为什么敢把生意交给它", 30, MAIN, True)])
bullets(s, 0.9, 1.5, 11.5, 4.6, [
    ("围栏：AI 的行动边界是规则，不是自觉", "报价必审、券库存熔断、线索出域必审、客资隐私脱敏——26 条酒店围栏事前裁决"),
    ("人审：关键决策永远有人拍板", "必审事项 AI 只提案不执行；审批卡片直达老板手机，一键批准/驳回"),
    ("可追溯：每个动作都有案底", "五元事件 + hash 链，谁、在什么时间、做了什么、依据哪条规则——可验链、可回放"),
    ("不越权：模型失灵也不失控", "离线确定性模型兜底；异常即挂起转人工，绝不静默硬闯"),
], size=17)

# ⑨ 业务↔代码对照
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("给开发者的桥：每个业务能力背后是哪块代码", 28, MAIN, True)])
rows = [
    ("获客五环 · 意图洞察", "social-listening + bundles/geo-growth 情报管线"),
    ("内容生产与分发", "ai-video + video-studio + publish-rpa（六平台 RPA）"),
    ("询盘承接与转化", "service-* 四件（C 端）+ im-channels + bundles/hotel 获客域"),
    ("AI 班组调度", "captain（ASK/QUEST 编排）+ night-shift（夜班）+ model-router"),
    ("围栏与人审", "fence-engine（DSL 裁决）+ review-console（人审台）"),
    ("可追溯与隔离", "workdata（五元事件/RLS）+ packages/db（迁移/验链）"),
]
for i, (a, b) in enumerate(rows):
    y = 1.45 + i * 0.95
    text(s, 0.9, y, 4.4, 0.9, [(a, 15, MAIN, True)])
    text(s, 5.5, y, 7.2, 0.9, [(b, 14, TEXT, False)])

# ⑩ 开箱体验
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("5 分钟开箱：不需要任何账号和密钥", 30, MAIN, True)])
text(s, 0.9, 1.5, 11.5, 1.0, [("pnpm install && pnpm preview:all", 24, MAIN, True, PP_ALIGN.CENTER)])
bullets(s, 0.9, 2.7, 11.5, 3.6, [
    ("PC 端 B 端工作台 → http://localhost:3000", "经营主页全员就位、晨报、待审批、一句话目标"),
    ("B 端移动 → http://localhost:3001", "12 页高保真演示页 + 手机壳容器"),
    ("C 端 AI 服务前台 → http://localhost:3002", "免登对话：查订单/售后/物流/常见问题"),
    ("Mock 数据已固化", "种子演示数据 + 离线确定性模型 + 演示直登；界面常驻「全模拟运行态」横幅"),
], size=16)

# ⑪ 二次开发路径
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("接手二次开发：人与 Agent 各有一条通道", 28, MAIN, True)])
bullets(s, 0.9, 1.5, 11.5, 4.8, [
    ("人类开发者", "README 能力速览 → install 欢迎横幅 → docs/capabilities.auto.md（能力导览）→ 本文 PPT"),
    ("AI Coding Agent", "AGENTS.md 入场 → .ai-prompt → pnpm agent:tour（能力自检巡游）→ docs/capability-map.md（全量机器清单）"),
    ("沙箱自带「手」", "computer-use 浏览器/电脑操作能力：preflight 后可用 browser_* 动作自动打开三端验证，UI 改动必须截图核对"),
    ("工程纪律（硬性）", "pnpm suite（445 条）· pnpm release:gate（发布门禁，未全过禁止发布）· docs/design-system.md（糖果视觉规范）"),
], size=16)

# ⑫ 行业复制
s = slide()
text(s, 0.8, 0.45, 11.7, 0.9, [("从酒店到千行百业：Bundle 化复制", 30, MAIN, True)])
bullets(s, 0.9, 1.5, 11.5, 4.4, [
    ("底座不变", "围栏/编排/事件/隔离/人审/夜班——行业无关的经营基建"),
    ("行业打包", "每个行业一个 Bundle：员工（preset）+ 围栏 + 技能 + 对象 + 管线 + 默认客群配置"),
    ("已有样例", "bundles/hotel（酒店获客 v3.3.0）· bundles/geo-growth（GEO 双域）· bundles/ai-video（短视频）"),
    ("四仓矩阵", "workloom（获客主仓）· workloom-im（IM 底座）· workloom-hotel（酒店垂直）· hyperreality-system（视频制作）"),
], size=17)

# ⑬ 结尾
s = slide(dark=True)
text(s, 1, 2.3, 11.3, 3, [
    ("让获客成为一门可测量的生意", 36, WHITE, True, PP_ALIGN.CENTER),
    ("pnpm preview:all —— 现在就用 5 分钟看全貌", 20, WHITE, False, PP_ALIGN.CENTER),
    ("文档索引：README · docs/SYSTEM-OVERVIEW.md（Agent 版详解）· docs/capabilities.auto.md（能力导览）", 14, RGBColor(0xFF, 0xD6, 0xDE), False, PP_ALIGN.CENTER),
])

out = os.path.join(ROOT, "docs/system-overview.pptx")
prs.save(out)
print(f"✓ {out}（{len(prs.slides._sldIdLst)} 页）")
